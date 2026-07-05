"""多格式文件加载器 — TXT/MD/DOCX/PPTX。"""
from pathlib import Path
from app.config.loader import get_config
from app.utils.log_tool import get_logger

logger = get_logger(__name__)


def _get_allow_types() -> set[str]:
    return set(get_config("allow_knowledge_file_types", ["txt", "pdf", "md", "pptx", "docx"]))


def _get_encodings() -> list[str]:
    return get_config("text_encodings", ["utf-8", "gbk", "gb2312", "latin-1"])


def load_file(file_path: str | Path, extension: str,
              user_id: str = "", md5_hex: str = "") -> list:
    file_path = Path(file_path)
    ext = extension.lower().lstrip(".")

    if ext == "txt":
        return txt_loader(file_path)
    elif ext == "md":
        return markdown_loader(file_path)
    elif ext == "docx":
        return docx_loader(file_path, user_id=user_id, md5_hex=md5_hex)
    elif ext == "pptx":
        return pptx_loader(file_path)
    else:
        raise ValueError(f"不支持的文件格式: .{ext}")


def txt_loader(file_path: Path) -> list:
    from langchain_community.document_loaders import TextLoader

    for encoding in _get_encodings():
        try:
            loader = TextLoader(str(file_path), encoding=encoding)
            docs = loader.load()
            if docs and docs[0].page_content.strip():
                logger.debug(f"【文本文件加载】使用编码 {encoding} 成功加载")
                return docs
        except Exception as e:
            logger.debug(f"【文本文件加载】编码 {encoding} 加载失败: {e}")

    logger.error("【文本文件加载】所有编码均失败")
    return []


def markdown_loader(file_path: Path) -> list:
    """Markdown 加载 — mistune v3 结构化解析 (TOC/表格/列表) → 降级 TextLoader。

    主路径: mistune v3 AST 解析，提取完整章节层级、表格、列表等结构。
    降级路径: 自定义正则剥离 MD 标记 → 按 TXT 逻辑 TextLoader 解析。
    """
    # 1. 读文件（编码降级链）
    content = None
    used_encoding = None
    for encoding in _get_encodings():
        try:
            content = file_path.read_text(encoding=encoding)
            if content.strip():
                used_encoding = encoding
                break
        except Exception as e:
            logger.debug(f"【Markdown文件加载】编码 {encoding} 读取失败: {e}")

    if content is None:
        logger.error("【Markdown文件加载】所有编码均读取失败")
        return []

    source = str(file_path)

    # 2. 主路径：mistune v3 结构化解析
    try:
        from app.utils.md_parser import parse_markdown
        structured_text, toc = parse_markdown(content)
        if structured_text.strip():
            from langchain_core.documents import Document as LCDoc
            import json
            logger.debug(f"【Markdown文件加载】mistune v3 解析成功 (编码: {used_encoding}, TOC: {len(toc)} 条)")
            current_chapter = toc[0]["text"] if toc else ""
            chapter_level = toc[0]["level"] if toc else 0
            return [LCDoc(
                page_content=structured_text,
                metadata={
                    "source": source,
                    "toc": json.dumps(toc, ensure_ascii=False),
                    "chapter_count": len(toc),
                    "current_chapter": current_chapter,
                    "chapter_level": chapter_level,
                },
            )]
    except Exception as e:
        logger.warning(f"【Markdown文件加载】mistune v3 解析失败: {e}，降级使用 TextLoader + 正则剥离")

    # 3. 降级路径：正则剥离 MD 标记 → TextLoader 解析
    from app.utils.md_parser import _strip_markdown_markup
    from langchain_community.document_loaders import TextLoader

    stripped = _strip_markdown_markup(content)
    if not stripped.strip():
        logger.error("【Markdown文件加载】正则剥离后内容为空")
        return []

    # 写入临时文件供 TextLoader 读取
    import tempfile
    import os
    with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as tf:
        tf.write(stripped)
        tmp_path = tf.name

    try:
        loader = TextLoader(tmp_path, encoding='utf-8')
        docs = loader.load()
        if docs and docs[0].page_content.strip():
            docs[0].metadata["source"] = source
            logger.debug(f"【Markdown文件加载】降级路径成功 (编码: {used_encoding})")
            return docs
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass

    logger.error("【Markdown文件加载】降级路径 TextLoader 返回空")
    return []


def _extract_docx_chapter(file_path: Path) -> tuple[str, int]:
    """从 DOCX 段落样式中提取第一个 Heading 作为章节信息。"""
    try:
        from docx import Document
        doc = Document(str(file_path))
        for p in doc.paragraphs:
            style = p.style.name if p.style else ""
            if style.startswith("Heading"):
                text = p.text.strip()
                if text:
                    try:
                        level = int(style.split()[-1])
                    except (ValueError, IndexError):
                        level = 0
                    return text, level
    except Exception:
        pass
    return "", 0


def _extract_docx_toc(file_path: Path) -> tuple[list[dict], int]:
    """从 DOCX 提取所有 Heading 作为 TOC 结构，返回 (toc_list, chapter_count)。"""
    toc = []
    chapter_path = []
    try:
        from docx import Document
        doc = Document(str(file_path))
        for p in doc.paragraphs:
            style = p.style.name if p.style else ""
            if style.startswith("Heading"):
                text = p.text.strip()
                if not text:
                    continue
                try:
                    level = int(style.split()[-1])
                except (ValueError, IndexError):
                    level = 0
                while len(chapter_path) >= level:
                    chapter_path.pop()
                chapter_path.append(text)
                toc.append({
                    "level": level,
                    "text": text,
                    "path": " > ".join(chapter_path),
                })
    except Exception:
        pass
    return toc, len(toc)


def _extract_docx_images(file_path: Path, user_id: str, md5: str) -> list[str]:
    """从 DOCX 提取内嵌图片，持久化到 extracted_images/，返回相对路径列表。"""
    image_paths = []
    try:
        from docx import Document
        from docx.opc.constants import RELATIONSHIP_TYPE as RT
        from app.utils.path_tool import get_data_path

        doc = Document(str(file_path))
        output_dir = get_data_path(f"extracted_images/{user_id}/{md5}")
        output_dir.mkdir(parents=True, exist_ok=True)

        for i, rel in enumerate(doc.part.rels.values()):
            if "image" not in rel.reltype:
                continue
            try:
                image_bytes = rel.target_part.blob
                ext = rel.target_part.partname.split(".")[-1] if "." in rel.target_part.partname else "png"
                img_filename = f"docx_i{i}.{ext}"
                img_path = output_dir / img_filename
                img_path.write_bytes(image_bytes)
                relative_path = img_path.relative_to(get_data_path()).as_posix()
                image_paths.append(relative_path)
            except Exception:
                continue
        if image_paths:
            logger.info(f"【DOCX图片提取】共提取 {len(image_paths)} 张图片")
    except ImportError:
        pass
    except Exception as e:
        logger.warning(f"【DOCX图片提取】失败: {e}")
    return image_paths


def docx_loader(file_path: Path, user_id: str = "", md5_hex: str = "") -> list:
    import json
    current_chapter, chapter_level = "", 0
    toc_list, chapter_count = [], 0
    image_paths = []

    try:
        current_chapter, chapter_level = _extract_docx_chapter(file_path)
        toc_list, chapter_count = _extract_docx_toc(file_path)
        if user_id and md5_hex:
            image_paths = _extract_docx_images(file_path, user_id, md5_hex)
    except Exception:
        pass

    has_images = len(image_paths) > 0
    base_meta = {
        "page": 1,
        "has_images": has_images,
        "toc": json.dumps(toc_list, ensure_ascii=False),
        "chapter_count": chapter_count,
        "current_chapter": current_chapter,
        "chapter_level": chapter_level,
    }
    if image_paths:
        base_meta["image_paths"] = image_paths

    # 主路径：Docx2txtLoader
    try:
        from langchain_community.document_loaders import Docx2txtLoader
        loader = Docx2txtLoader(str(file_path))
        docs = loader.load()
        if docs and docs[0].page_content.strip():
            docs[0].metadata.update(base_meta)
            docs[0].metadata["source"] = str(file_path)
            logger.debug(f"【WORD文件加载】Docx2txtLoader 成功 (TOC: {chapter_count}, 图片: {len(image_paths)})")
            return docs
    except Exception as e:
        logger.warning(f"【WORD文件加载】Docx2txtLoader 失败: {e}，尝试 python-docx 兜底")

    # 降级路径：python-docx
    try:
        from docx import Document
        from langchain_core.documents import Document as LCDoc
        doc = Document(str(file_path))
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        if text.strip():
            logger.debug("【WORD文件加载】python-docx 兜底成功")
            return [LCDoc(page_content=text, metadata={**base_meta, "source": str(file_path)})]
    except Exception as e:
        logger.error(f"【WORD文件加载】python-docx 兜底也失败: {e}")

    return []


def pptx_loader(file_path: Path) -> list:
    """PPTX 加载 — python-pptx 原生解析。解析失败直接抛异常，无降级。"""
    from pptx import Presentation
    from langchain_core.documents import Document as LCDoc

    prs = Presentation(str(file_path))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
    text = "\n".join(texts)
    if not text.strip():
        raise ValueError(f"PPTX 文件解析结果为空: {file_path.name}")
    logger.debug(f"【PPT文件加载】python-pptx 成功加载: {len(texts)} 个文本片段")
    return [LCDoc(page_content=text, metadata={"source": str(file_path)})]
